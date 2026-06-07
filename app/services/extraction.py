import csv
import email
import hashlib
import mimetypes
import uuid
from dataclasses import dataclass
from datetime import datetime
from email.header import decode_header, make_header
from email.policy import default
from email.utils import parseaddr, parsedate_to_datetime
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text
from werkzeug.utils import secure_filename

from app.extensions import db
from app.models import (
    DocumentChunk,
    DocumentText,
    EmailMessage,
    EmailRecipient,
    ProcessingLog,
    SourceFile,
    utcnow,
)


@dataclass
class ExtractionResult:
    text: str
    method: str
    metadata: dict | None = None


def _now():
    return utcnow()


def _safe_decode(value):
    if not value:
        return None

    try:
        return str(make_header(decode_header(value)))
    except Exception:
        return str(value)


def _parse_email_address(value):
    name, addr = parseaddr(value or "")
    return _safe_decode(name), addr


def _parse_email_date(value):
    if not value:
        return None

    try:
        return parsedate_to_datetime(value)
    except Exception:
        return None


def _log(source_file_id, stage, status, message=None, started_at=None):
    finished_at = _now()
    duration_ms = None

    if started_at:
        duration_ms = int((finished_at - started_at).total_seconds() * 1000)

    entry = ProcessingLog(
        source_file_id=source_file_id,
        stage=stage,
        status=status,
        message=message,
        started_at=started_at,
        finished_at=finished_at,
        duration_ms=duration_ms,
    )
    db.session.add(entry)


def _calculate_sha256(file_path: Path) -> str:
    sha = hashlib.sha256()

    with file_path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            sha.update(block)

    return sha.hexdigest()


def _estimate_tokens(text: str) -> int:
    if not text:
        return 0

    return max(1, int(len(text) / 4))


def _normalise_text(text: str) -> str:
    if not text:
        return ""

    text = text.replace("\x00", "")
    lines = [line.rstrip() for line in text.splitlines()]
    compact = "\n".join(lines)
    while "\n\n\n" in compact:
        compact = compact.replace("\n\n\n", "\n\n")

    return compact.strip()


def _chunk_text(text: str, max_chars: int = 4500) -> list[str]:
    text = _normalise_text(text)

    if not text:
        return []

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    current = ""

    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= max_chars:
            current = f"{current}\n\n{paragraph}".strip()
        else:
            if current:
                chunks.append(current)
            current = paragraph

    if current:
        chunks.append(current)

    if not chunks and text:
        chunks = [text[i:i + max_chars] for i in range(0, len(text), max_chars)]

    return chunks


def extract_txt_like(file_path: Path) -> ExtractionResult:
    text = file_path.read_text(encoding="utf-8", errors="ignore")
    return ExtractionResult(text=text, method="plain_text")


def extract_html(file_path: Path) -> ExtractionResult:
    html = file_path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n")
    return ExtractionResult(text=text, method="html_bs4")


def extract_rtf(file_path: Path) -> ExtractionResult:
    raw = file_path.read_text(encoding="utf-8", errors="ignore")
    return ExtractionResult(text=rtf_to_text(raw), method="striprtf")


def extract_csv(file_path: Path) -> ExtractionResult:
    rows = []

    with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
        sample = f.read(4096)
        f.seek(0)

        try:
            dialect = csv.Sniffer().sniff(sample)
        except Exception:
            dialect = csv.excel

        reader = csv.reader(f, dialect)

        for row_index, row in enumerate(reader, start=1):
            cleaned = [cell.strip() for cell in row]
            rows.append(f"Row {row_index}: " + " | ".join(cleaned))

    return ExtractionResult(text="\n".join(rows), method="csv_reader")


def extract_pdf(file_path: Path) -> ExtractionResult:
    reader = PdfReader(str(file_path))
    parts = []

    for index, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception:
            page_text = ""

        if page_text.strip():
            parts.append(f"--- Page {index} ---\n{page_text.strip()}")

    return ExtractionResult(
        text="\n\n".join(parts),
        method="pypdf",
        metadata={"page_count": len(reader.pages)},
    )


def extract_docx(file_path: Path) -> ExtractionResult:
    document = DocxDocument(str(file_path))
    parts = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        parts.append(f"\n--- Table {table_index} ---")
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            parts.append(" | ".join(cells))

    return ExtractionResult(text="\n".join(parts), method="python_docx")


def extract_xlsx(file_path: Path) -> ExtractionResult:
    workbook = load_workbook(filename=str(file_path), data_only=True, read_only=True)
    parts = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")

        for row in sheet.iter_rows(values_only=True):
            values = []
            for value in row:
                if value is None:
                    values.append("")
                elif isinstance(value, datetime):
                    values.append(value.isoformat())
                else:
                    values.append(str(value))

            if any(cell.strip() for cell in values):
                parts.append(" | ".join(values))

        parts.append("")

    return ExtractionResult(
        text="\n".join(parts),
        method="openpyxl",
        metadata={"sheet_names": workbook.sheetnames},
    )


def extract_pptx(file_path: Path) -> ExtractionResult:
    presentation = Presentation(str(file_path))
    parts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"--- Slide {slide_index} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)

        if slide.has_notes_slide:
            notes = []
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    notes.append(shape.text.strip())

            if notes:
                parts.append("Speaker notes:")
                parts.extend(notes)

        parts.append("")

    return ExtractionResult(
        text="\n".join(parts),
        method="python_pptx",
        metadata={"slide_count": len(presentation.slides)},
    )


def _payload_to_text(payload) -> str:
    if payload is None:
        return ""

    charset = payload.get_content_charset() or "utf-8"

    try:
        raw = payload.get_payload(decode=True)
        if raw:
            return raw.decode(charset, errors="ignore")
    except Exception:
        pass

    try:
        return str(payload.get_content())
    except Exception:
        return ""


def extract_eml(file_path: Path, source_file: SourceFile) -> ExtractionResult:
    with file_path.open("rb") as f:
        msg = email.message_from_binary_file(f, policy=default)

    subject = _safe_decode(msg.get("subject"))
    message_id = msg.get("message-id")
    sender_name, sender_email = _parse_email_address(msg.get("from"))
    sent_at = _parse_email_date(msg.get("date"))

    plain_parts = []
    html_parts = []

    attachment_count = 0

    for part in msg.walk():
        content_disposition = part.get_content_disposition()
        content_type = part.get_content_type()

        if content_disposition == "attachment":
            attachment_count += 1
            _save_eml_attachment(part, source_file)
            continue

        if content_type == "text/plain":
            plain_parts.append(_payload_to_text(part))
        elif content_type == "text/html":
            html = _payload_to_text(part)
            html_parts.append(html)

    body_text = "\n\n".join([p for p in plain_parts if p.strip()])

    if not body_text and html_parts:
        soup = BeautifulSoup("\n\n".join(html_parts), "html.parser")
        body_text = soup.get_text(separator="\n")

    email_record = EmailMessage.query.filter_by(source_file_id=source_file.id).first()

    if not email_record:
        email_record = EmailMessage(source_file_id=source_file.id)
        db.session.add(email_record)

    email_record.message_id = message_id
    email_record.subject = subject
    email_record.sender_name = sender_name
    email_record.sender_email = sender_email
    email_record.sent_at = sent_at
    email_record.body_text = body_text
    email_record.thread_key = subject

    EmailRecipient.query.filter_by(email_message_id=email_record.id).delete()

    db.session.flush()

    for header_name, recipient_type in [("to", "to"), ("cc", "cc"), ("bcc", "bcc")]:
        values = msg.get_all(header_name, [])
        for value in values:
            for item in str(value).split(","):
                name, addr = _parse_email_address(item)
                if name or addr:
                    db.session.add(
                        EmailRecipient(
                            email_message_id=email_record.id,
                            recipient_type=recipient_type,
                            name=name,
                            email=addr,
                        )
                    )

    text = "\n".join(
        [
            f"Subject: {subject or ''}",
            f"From: {sender_name or ''} <{sender_email or ''}>",
            f"Date: {sent_at.isoformat() if sent_at else ''}",
            "",
            body_text or "",
        ]
    )

    return ExtractionResult(
        text=text,
        method="eml_parser",
        metadata={"attachment_count": attachment_count},
    )


def _save_eml_attachment(part, parent_source_file: SourceFile):
    filename = part.get_filename()

    if not filename:
        filename = f"attachment-{uuid.uuid4().hex}"

    filename = _safe_decode(filename) or f"attachment-{uuid.uuid4().hex}"
    safe_name = secure_filename(filename)
    file_ext = Path(safe_name).suffix.lower()
    stored_filename = f"{uuid.uuid4().hex}{file_ext}"

    attachments_dir = Path(parent_source_file.storage_path).resolve().parents[1] / "attachments"
    attachments_dir.mkdir(parents=True, exist_ok=True)

    destination = attachments_dir / stored_filename

    payload = part.get_payload(decode=True) or b""
    destination.write_bytes(payload)

    file_size = destination.stat().st_size
    sha256_hash = _calculate_sha256(destination)
    mime_type, _ = mimetypes.guess_type(str(destination))

    existing = SourceFile.query.filter_by(
        sha256_hash=sha256_hash,
        file_size=file_size,
    ).first()

    if existing:
        destination.unlink(missing_ok=True)
        return

    child = SourceFile(
        original_filename=filename,
        stored_filename=stored_filename,
        file_ext=file_ext,
        mime_type=mime_type,
        file_size=file_size,
        sha256_hash=sha256_hash,
        storage_path=str(destination),
        source_type="email_attachment",
        upload_method="email_attachment",
        parent_file_id=parent_source_file.id,
        processing_status="uploaded",
        business_area=parent_source_file.business_area,
        uploaded_by_id=parent_source_file.uploaded_by_id,
    )

    db.session.add(child)
    db.session.flush()

    _log(
        child.id,
        "attachment_extract",
        "success",
        f"Extracted from parent email: {parent_source_file.original_filename}",
        _now(),
    )


def extract_msg(file_path: Path, source_file: SourceFile) -> ExtractionResult:
    import extract_msg

    msg = extract_msg.Message(str(file_path))
    msg_sender = msg.sender or ""
    sender_name, sender_email = _parse_email_address(msg_sender)

    body = msg.body or ""
    subject = msg.subject or source_file.original_filename

    sent_at = None
    try:
        sent_at = msg.date
    except Exception:
        sent_at = None

    email_record = EmailMessage.query.filter_by(source_file_id=source_file.id).first()

    if not email_record:
        email_record = EmailMessage(source_file_id=source_file.id)
        db.session.add(email_record)

    email_record.subject = subject
    email_record.sender_name = sender_name
    email_record.sender_email = sender_email
    email_record.sent_at = sent_at
    email_record.body_text = body
    email_record.thread_key = subject

    attachment_count = 0

    attachments_dir = Path(source_file.storage_path).resolve().parents[1] / "attachments"
    attachments_dir.mkdir(parents=True, exist_ok=True)

    for attachment in msg.attachments:
        attachment_count += 1

        filename = getattr(attachment, "longFilename", None) or getattr(attachment, "shortFilename", None)
        if not filename:
            filename = f"attachment-{uuid.uuid4().hex}"

        safe_name = secure_filename(filename)
        file_ext = Path(safe_name).suffix.lower()
        stored_filename = f"{uuid.uuid4().hex}{file_ext}"
        destination = attachments_dir / stored_filename

        try:
            data = attachment.data
            destination.write_bytes(data)
        except Exception:
            continue

        file_size = destination.stat().st_size
        sha256_hash = _calculate_sha256(destination)
        mime_type, _ = mimetypes.guess_type(str(destination))

        existing = SourceFile.query.filter_by(
            sha256_hash=sha256_hash,
            file_size=file_size,
        ).first()

        if existing:
            destination.unlink(missing_ok=True)
            continue

        child = SourceFile(
            original_filename=filename,
            stored_filename=stored_filename,
            file_ext=file_ext,
            mime_type=mime_type,
            file_size=file_size,
            sha256_hash=sha256_hash,
            storage_path=str(destination),
            source_type="email_attachment",
            upload_method="email_attachment",
            parent_file_id=source_file.id,
            processing_status="uploaded",
            business_area=source_file.business_area,
            uploaded_by_id=source_file.uploaded_by_id,
        )

        db.session.add(child)
        db.session.flush()

        _log(
            child.id,
            "attachment_extract",
            "success",
            f"Extracted from parent email: {source_file.original_filename}",
            _now(),
        )

    text = "\n".join(
        [
            f"Subject: {subject or ''}",
            f"From: {sender_name or ''} <{sender_email or ''}>",
            f"Date: {sent_at.isoformat() if sent_at else ''}",
            "",
            body or "",
        ]
    )

    return ExtractionResult(
        text=text,
        method="extract_msg",
        metadata={"attachment_count": attachment_count},
    )


def extract_source_file(source_file_id: int) -> SourceFile:
    source_file = SourceFile.query.get(source_file_id)

    if not source_file:
        raise ValueError(f"SourceFile not found: {source_file_id}")

    started_at = _now()
    source_file.processing_status = "extracting"
    source_file.processing_error = None
    db.session.commit()

    _log(source_file.id, "extraction", "started", "Text extraction started.", started_at)
    db.session.commit()

    file_path = Path(source_file.storage_path)

    try:
        if not file_path.exists():
            raise FileNotFoundError(f"Stored file missing: {file_path}")

        ext = (source_file.file_ext or file_path.suffix or "").lower()

        if ext in {".txt", ".md"}:
            result = extract_txt_like(file_path)
        elif ext == ".html":
            result = extract_html(file_path)
        elif ext == ".rtf":
            result = extract_rtf(file_path)
        elif ext == ".csv":
            result = extract_csv(file_path)
        elif ext == ".pdf":
            result = extract_pdf(file_path)
        elif ext == ".docx":
            result = extract_docx(file_path)
        elif ext == ".xlsx":
            result = extract_xlsx(file_path)
        elif ext == ".pptx":
            result = extract_pptx(file_path)
        elif ext == ".eml":
            result = extract_eml(file_path, source_file)
        elif ext == ".msg":
            result = extract_msg(file_path, source_file)
        else:
            raise ValueError(f"Unsupported extraction type for {ext}")

        text = _normalise_text(result.text)

        extracted_dir = file_path.resolve().parents[1] / "extracted_text"
        extracted_dir.mkdir(parents=True, exist_ok=True)

        text_filename = f"{source_file.id}_{Path(source_file.stored_filename).stem}.txt"
        text_path = extracted_dir / text_filename
        text_path.write_text(text, encoding="utf-8", errors="ignore")

        existing_text = DocumentText.query.filter_by(source_file_id=source_file.id).first()

        if not existing_text:
            existing_text = DocumentText(source_file_id=source_file.id)
            db.session.add(existing_text)

        existing_text.extracted_text_path = str(text_path)
        existing_text.text_preview = text[:4000]
        existing_text.word_count = len(text.split())
        existing_text.char_count = len(text)
        existing_text.extraction_method = result.method
        existing_text.extraction_status = "success"

        DocumentChunk.query.filter_by(source_file_id=source_file.id).delete()

        chunks = _chunk_text(text)

        for index, chunk in enumerate(chunks, start=1):
            db.session.add(
                DocumentChunk(
                    source_file_id=source_file.id,
                    chunk_index=index,
                    chunk_text=chunk,
                    token_estimate=_estimate_tokens(chunk),
                )
            )

        source_file.processing_status = "extracted"
        source_file.processed_at = _now()
        source_file.processing_error = None

        _log(
            source_file.id,
            "extraction",
            "success",
            f"Extracted {len(text)} characters into {len(chunks)} chunk(s) using {result.method}.",
            started_at,
        )

        db.session.commit()
        return source_file

    except Exception as exc:
        source_file.processing_status = "failed"
        source_file.processing_error = str(exc)

        existing_text = DocumentText.query.filter_by(source_file_id=source_file.id).first()

        if existing_text:
            existing_text.extraction_status = "failed"

        _log(source_file.id, "extraction", "failed", str(exc), started_at)
        db.session.commit()
        raise
